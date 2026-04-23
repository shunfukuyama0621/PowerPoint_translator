"""PowerPoint (.pptx) の翻訳。書式・画像・レイアウトを維持する。"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from translator import translate_texts


def _collect_runs(shape, runs: list):
    """図形を再帰的に走査してテキストrunを集める（グループ・表も対応）。"""
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub in shape.shapes:
            _collect_runs(sub, runs)
        return

    # テーブル
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            runs.append(run)
        return

    # 通常のテキストフレーム
    if shape.has_text_frame:
        # 枠内に文字が収まらなくなった場合に自動縮小する
        try:
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            shape.text_frame.word_wrap = True
        except Exception:
            pass

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    runs.append(run)


def translate_pptx(
    input_bytes: bytes,
    target_lang: str,
    source_lang: str = "ja",
) -> bytes:
    """PPTXバイト列を受け取り、翻訳済みのPPTXバイト列を返す。

    Args:
        target_lang / source_lang: UI言語コード（"ja" / "en" / "zh"）
    """
    prs = Presentation(BytesIO(input_bytes))

    all_runs: list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            _collect_runs(shape, all_runs)

    if not all_runs:
        out = BytesIO()
        prs.save(out)
        return out.getvalue()

    original_texts = [r.text for r in all_runs]
    translated = translate_texts(original_texts, target_lang, source_lang)

    for run, original_text, new_text in zip(all_runs, original_texts, translated):
        # 訳文が大きく伸びた場合、フォントサイズを段階的に縮小
        # （通常は auto_size が処理するが保険として）
        if (
            run.font.size is not None
            and len(original_text) > 0
            and len(new_text) > len(original_text) * 1.5
        ):
            current_pt = run.font.size.pt
            if current_pt > 10:
                run.font.size = Pt(max(current_pt * 0.9, 9))

        run.text = new_text

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
